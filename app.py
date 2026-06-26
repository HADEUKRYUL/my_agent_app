import streamlit as st
from openai import OpenAI
import PyPDF2
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import base64
import json
import os
import datetime
import math
import urllib.request
from io import BytesIO
from apscheduler.schedulers.background import BackgroundScheduler

# 🌐 언어별 레이블 (다국어 UI 지원)
LABELS = {
    "한국어": {
        "title": "🤖 나만의 특급 비서 에이전트",
        "tabs": ["💬 대화", "📅 일정", "📊 생산실적", "📚 매뉴얼(자가학습)", "💼 코어워크(PPT)", "📈 데이터분석"],
        "sb_chat": "💬 대화 세션 관리", "sb_new": "➕ 새 대화 시작", "sb_search": "🔍 대화 검색",
        "sb_upload": "📁 실적/데이터 업로드", "sb_trans": "🌐 퀵 번역기", "trans_btn": "🚀 번역 실행"
    },
    "베트남어": {
        "title": "🤖 Trợ lý AI Đặc biệt",
        "tabs": ["💬 Trò chuyện", "📅 Lịch trình", "📊 NS Sản xuất", "📚 Sổ tay", "💼 Báo cáo(PPT)", "📈 Phân tích Dữ liệu"],
        "sb_chat": "💬 Quản lý phiên", "sb_new": "➕ Trò chuyện mới", "sb_search": "🔍 Tìm kiếm",
        "sb_upload": "📁 Tải lên dữ liệu", "sb_trans": "🌐 Dịch nhanh", "trans_btn": "🚀 Dịch"
    },
    "영어": {
        "title": "🤖 Special AI Assistant",
        "tabs": ["💬 Chat", "📅 Schedule", "📊 Production", "📚 Manual", "💼 Core Work(PPT)", "📈 Data Analysis"],
        "sb_chat": "💬 Session Management", "sb_new": "➕ New Chat", "sb_search": "🔍 Search",
        "sb_upload": "📁 Upload Data", "sb_trans": "🌐 Quick Translator", "trans_btn": "🚀 Translate"
    }
}

# 💾 1. 영구 저장소
HISTORY_FILE, SCHEDULE_FILE, PRODUCTION_FILE = "chat_history.json", "schedule_history.json", "production_history.json"
MANUAL_KNOWLEDGE_FILE, MANUAL_CHAT_FILE = "manual_knowledge.json", "manual_chat_history.json"

def load_json(file_path, default_type=dict):
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
                if default_type == dict and isinstance(data, list):
                    return {"current_id": "default", "sessions": [{"id": "default", "title": "기본 대화방", "messages": data}]}
                return data
        except: return default_type()
    return default_type() if default_type == dict else []

def save_json(file_path, data):
    with open(file_path, "w", encoding="utf-8") as f: json.dump(data, f, ensure_ascii=False, indent=4)

def safe_float(val):
    try:
        if pd.isna(val): return 0.0
        val_str = str(val).replace(',', '').strip()
        if not val_str or val_str.lower() in ['nan', 'nat', 'null']: return 0.0
        f = float(val_str)
        return 0.0 if math.isnan(f) else f
    except: return 0.0

if "scheduler" not in st.session_state:
    st.session_state.scheduler = BackgroundScheduler()
    st.session_state.scheduler.start()

st.set_page_config(page_title="AI Assistant", page_icon="🤖", layout="wide")

# 🌍 다국어 선택기
if "lang" not in st.session_state: st.session_state.lang = "한국어"
col_t1, col_t2 = st.columns([8, 2])
with col_t2: st.session_state.lang = st.selectbox("Language / Ngôn ngữ", ["한국어", "베트남어", "영어"])
L = LABELS[st.session_state.lang]

st.title(L["title"])

# API 세팅
try: my_key = st.secrets["OPENAI_API_KEY"]
except: my_key = None
if not my_key:
    st.error("🚨 설정(Secrets)에서 OPENAI_API_KEY를 입력하세요.")
    st.stop()
client = OpenAI(api_key=my_key)

# 세션 초기화
if "chats" not in st.session_state:
    st.session_state.chats = load_json(HISTORY_FILE, default_type=dict)
    if not st.session_state.chats or "sessions" not in st.session_state.chats:
        st.session_state.chats = {"current_id": "default", "sessions": [{"id": "default", "title": "기본 대화방", "messages": []}]}
if "schedules" not in st.session_state: st.session_state.schedules = load_json(SCHEDULE_FILE, default_type=list)
if "production_history" not in st.session_state: st.session_state.production_history = load_json(PRODUCTION_FILE, default_type=list)
if "manual_knowledge" not in st.session_state: st.session_state.manual_knowledge = load_json(MANUAL_KNOWLEDGE_FILE, default_type=list)
if "manual_chat" not in st.session_state: st.session_state.manual_chat = load_json(MANUAL_CHAT_FILE, default_type=list)
if "data_chat" not in st.session_state: st.session_state.data_chat = []
if "trans_res" not in st.session_state: st.session_state.trans_res = ""

# ==========================================
# 📁 사이드바
# ==========================================
with st.sidebar:
    st.header(L["sb_chat"])
    if st.button(L["sb_new"], use_container_width=True):
        new_id = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
        st.session_state.chats["sessions"].append({"id": new_id, "title": f"New Chat {datetime.datetime.now().strftime('%m-%d %H:%M')}", "messages": []})
        st.session_state.chats["current_id"] = new_id
        save_json(HISTORY_FILE, st.session_state.chats)
        st.rerun()

    search_query = st.text_input(L["sb_search"], "").strip()
    valid_sessions = [s for s in st.session_state.chats["sessions"] if not search_query or search_query.lower() in s["title"].lower() or any(search_query.lower() in str(m.get("content", "")).lower() for m in s.get("messages", []))]
    if valid_sessions:
        session_dict = {s["id"]: s["title"] for s in valid_sessions}
        curr_id = st.session_state.chats["current_id"]
        if curr_id not in session_dict: curr_id = valid_sessions[0]["id"]
        st.session_state.chats["current_id"] = st.selectbox("Select Session", list(session_dict.keys()), format_func=lambda x: session_dict[x], index=list(session_dict.keys()).index(curr_id), label_visibility="collapsed")

    st.markdown("---")
    st.header(L["sb_upload"])
    uploaded_file = st.file_uploader("", type=["pdf", "csv", "xlsx", "txt", "pptx", "png", "jpg", "jpeg"])
    file_content, image_data, mime_type = None, None, None

    if uploaded_file:
        try:
            if uploaded_file.type.startswith("image"):
                image_data = base64.b64encode(uploaded_file.read()).decode('utf-8')
                mime_type = uploaded_file.type
                st.image(uploaded_file)
                st.success("✅ Image Read!")
            elif uploaded_file.name.endswith('.csv') or uploaded_file.name.endswith('.xlsx'):
                df_upload = pd.read_csv(uploaded_file) if uploaded_file.name.endswith('.csv') else pd.read_excel(uploaded_file)
                file_content = df_upload.to_string()
                
                workplace_col, target_col, passed_col, defect_col = None, None, None, None
                for col in df_upload.columns:
                    col_str = str(col).lower().strip()
                    if any(k in col_str for k in ['작업장', 'workplace', '라인', 'line']): workplace_col = col
                    elif any(k in col_str for k in ['목표', 'target']): target_col = col
                    elif any(k in col_str for k in ['양품', 'passed', '합격']): passed_col = col
                    elif any(k in col_str for k in ['불량', 'defect']): defect_col = col

                if workplace_col and target_col and passed_col and defect_col:
                    today_str = datetime.date.today().strftime("%Y-%m-%d")
                    for _, row in df_upload.iterrows():
                        wp_name = str(row[workplace_col]).strip()
                        if wp_name.lower() in ['nan', 'nat', 'null', ''] or pd.isna(row[workplace_col]): continue
                        st.session_state.production_history.append({"Date": today_str, "Workplace": wp_name, "Target": safe_float(row[target_col]), "Passed": safe_float(row[passed_col]), "Defect": safe_float(row[defect_col])})
                    save_json(PRODUCTION_FILE, st.session_state.production_history)
                    st.success("✅ DB Saved!")
                else: st.warning("⚠️ Missing columns")
            elif uploaded_file.name.endswith('.pdf'): file_content = "\n".join([p.extract_text() for p in PyPDF2.PdfReader(uploaded_file).pages])
            elif uploaded_file.name.endswith('.pptx'): file_content = "\n".join([shape.text for slide in Presentation(uploaded_file).slides for shape in slide.shapes if hasattr(shape, "text")])
            else: file_content = uploaded_file.read().decode("utf-8")
        except Exception as e: st.error(f"Error: {e}")

    st.markdown("---")
    st.header(L["sb_trans"])
    target_lang = st.radio("To:", ["베트남어", "영어", "한국어"], horizontal=True, label_visibility="collapsed")
    translate_text = st.text_area("Text:", height=100, label_visibility="collapsed")
    if st.button(L["trans_btn"], use_container_width=True):
        if translate_text.strip():
            with st.spinner("Translating..."):
                try:
                    res = client.chat.completions.create(model="gpt-4o", messages=[{"role": "system", "content": f"Translate naturally to {target_lang}. Only output the result."}, {"role": "user", "content": translate_text}], temperature=0.3)
                    st.session_state.trans_res = res.choices[0].message.content
                except: pass
    if st.session_state.trans_res:
        st.info(st.session_state.trans_res)

# ==========================================
# 📱 메인 탭 구성
# ==========================================
tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(L["tabs"])
current_session = next(s for s in st.session_state.chats["sessions"] if s["id"] == st.session_state.chats["current_id"])

# 탭 1: 대화
with tab1:
    with st.expander(f"📝 {current_session['title']}"):
        new_title = st.text_input("Title", current_session["title"], label_visibility="collapsed")
        if st.button("Save"):
            current_session["title"] = new_title
            save_json(HISTORY_FILE, st.session_state.chats)
            st.rerun()

    for msg in current_session["messages"]:
        if msg["role"] != "system":
            with st.chat_message(msg["role"]): st.markdown(msg["content"])

    if prompt := st.chat_input("메시지 입력 (/그리기 내용)", key="main_chat"):
        current_session["messages"].append({"role": "user", "content": prompt})
        with st.chat_message("user"): st.markdown(prompt)

        if prompt.startswith("/그리기"):
            with st.chat_message("assistant"):
                with st.spinner("Drawing..."):
                    try:
                        draw_prompt = prompt.replace("/그리기", "").strip()
                        response = client.images.generate(model="dall-e-3", prompt=draw_prompt, size="1024x1024")
                        st.image(response.data[0].url)
                        current_session["messages"].append({"role": "assistant", "content": f"일러스트 완성: {draw_prompt}"})
                        save_json(HISTORY_FILE, st.session_state.chats)
                    except Exception as e: st.error(e)
        else:
            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    sys_inst = "당신은 주인님의 개인 비서입니다. 베트남 현지 지식과 데이터를 바탕으로 대답하세요."
                    msgs = [{"role": "system", "content": sys_inst}] + current_session["messages"][:-1]
                    if image_data: msgs.append({"role": "user", "content": [{"type": "text", "text": prompt}, {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_data}"}}]})
                    elif file_content: msgs.append({"role": "user", "content": f"{prompt}\n\n[Data]:\n{file_content[:5000]}"})
                    else: msgs.append({"role": "user", "content": prompt})
                    try:
                        res = client.chat.completions.create(model="gpt-4o", messages=msgs, temperature=0.7)
                        reply = res.choices[0].message.content
                        st.markdown(reply)
                        current_session["messages"].append({"role": "assistant", "content": reply})
                        save_json(HISTORY_FILE, st.session_state.chats)
                    except Exception as e: st.error(e)

# 탭 2: 일정
with tab2:
    with st.expander("➕ 새 일정 등록", expanded=True):
        c1, c2 = st.columns(2)
        with c1: nd = st.date_input("Date")
        with c1: nt = st.time_input("Time")
        with c2: n_task = st.text_input("Task")
        if st.button("Save Task"):
            if n_task:
                st.session_state.schedules.append({"Task": n_task, "DateTime": f"{nd} {nt.strftime('%H:%M')}"})
                save_json(SCHEDULE_FILE, st.session_state.schedules)
                st.rerun()
    if st.session_state.schedules:
        st.dataframe(pd.DataFrame(st.session_state.schedules).sort_values(by="DateTime").reset_index(drop=True), use_container_width=True)

# 탭 3: 생산실적
with tab3:
    if st.session_state.production_history:
        df_hist = pd.DataFrame(st.session_state.production_history)
        if 'Workplace' not in df_hist.columns and 'Line' in df_hist.columns: df_hist['Workplace'] = df_hist['Line']
        df_hist['Date_dt'] = pd.to_datetime(df_hist['Date'])
        curr_mo = datetime.date.today().strftime('%Y-%m')
        st.markdown(f"### 📅 당월 ({curr_mo}) 작업장별 현황")
        df_curr = df_hist[df_hist['Date_dt'].dt.strftime('%Y-%m') == curr_mo]
        if not df_curr.empty:
            df_grp = df_curr.groupby('Workplace').agg({'Target': 'sum', 'Passed': 'sum', 'Defect': 'sum'}).reset_index()
            df_grp['가동율'] = df_grp.apply(lambda r: r['Passed']/r['Target'] if r['Target']>0 else 0.0, axis=1)
            df_grp['PPM'] = df_grp.apply(lambda r: (r['Defect']/r['Passed'])*1000000 if r['Passed']>0 else 0.0, axis=1)
            
            st.dataframe(df_grp.rename(columns={'Workplace':'작업장명', 'Target':'목표', 'Passed':'양품', 'Defect':'불량'}).style.format({'목표': '{:,.0f}', '양품': '{:,.0f}', '불량': '{:,.0f}', '가동율': '{:.1%}', 'PPM': '{:,.0f} PPM'}), use_container_width=True)
            
            c1, c2 = st.columns(2)
            with c1:
                st.markdown("**📉 가동율 <= 80%**")
                df_u80 = df_grp[df_grp['가동율'] <= 0.8]
                if not df_u80.empty: st.bar_chart(df_u80.set_index('Workplace')[['가동율']])
            with c2:
                st.markdown("**🔥 PPM >= 500**")
                df_o500 = df_grp[df_grp['PPM'] >= 500]
                if not df_o500.empty: st.bar_chart(df_o500.set_index('Workplace')[['PPM']])
        
        st.markdown("### 📈 일자별 누적")
        df_daily = df_hist.groupby(df_hist['Date_dt'].dt.strftime('%Y-%m-%d')).agg({'Target':'sum', 'Passed':'sum'}).sort_index()
        df_daily['누적 목표'] = df_daily['Target'].cumsum()
        df_daily['누적 양품'] = df_daily['Passed'].cumsum()
        st.line_chart(df_daily[['누적 목표', '누적 양품']])
    else: st.info("No Data")

# 탭 4: 매뉴얼 챗봇
with tab4:
    with st.expander("🛠️ 매뉴얼 관리"):
        m_file = st.file_uploader("Upload Manual", type=["pdf", "txt", "pptx"])
        if st.button("🧠 훈련 실행") and m_file:
            with st.spinner("Learning..."):
                t = ""
                if m_file.name.endswith('.pdf'): t = "\n".join([p.extract_text() for p in PyPDF2.PdfReader(m_file).pages])
                elif m_file.name.endswith('.pptx'): t = "\n".join([s.text for sl in Presentation(m_file).slides for s in sl.shapes if hasattr(s, "text")])
                else: t = m_file.read().decode("utf-8")
                st.session_state.manual_knowledge.append({"title": m_file.name, "content": t[:30000]})
                save_json(MANUAL_KNOWLEDGE_FILE, st.session_state.manual_knowledge)
                st.rerun()
        if st.session_state.manual_knowledge:
            for i, doc in enumerate(st.session_state.manual_knowledge):
                c1, c2 = st.columns([8, 2])
                c1.markdown(f"📄 {doc['title']}")
                if c2.button("🗑️ 삭제", key=f"del_{i}"):
                    st.session_state.manual_knowledge.pop(i)
                    save_json(MANUAL_KNOWLEDGE_FILE, st.session_state.manual_knowledge)
                    st.rerun()

    st.markdown("---")
    for msg in st.session_state.manual_chat:
        if msg["role"] != "system":
            with st.chat_message(msg["role"]): st.markdown(msg["content"])

    if mp := st.chat_input("질문 (베트남어 번역 요청 가능)", key="m_chat"):
        st.session_state.manual_chat.append({"role": "user", "content": mp})
        with st.chat_message("user"): st.markdown(mp)
        with st.chat_message("assistant"):
            with st.spinner("Scanning..."):
                ctx = "\n".join([f"[{d['title']}]\n{d['content']}" for d in st.session_state.manual_knowledge])
                sys_m = "당신은 매뉴얼 전문가이자 번역가입니다. 제공된 매뉴얼 팩트 기반으로 대답하세요. '베트남어로 번역해' 등의 요구가 있으면 현지인이 완벽히 이해하게 베트남어로 번역하세요.\n" + ctx
                try:
                    res = client.chat.completions.create(model="gpt-4o", messages=[{"role": "system", "content": sys_m}] + st.session_state.manual_chat[:-1] + [{"role": "user", "content": mp}], temperature=0.3)
                    reply = res.choices[0].message.content
                    st.markdown(reply)
                    st.session_state.manual_chat.append({"role": "assistant", "content": reply})
                    save_json(MANUAL_CHAT_FILE, st.session_state.manual_chat)
                    
                    auto_doc = next((d for d in st.session_state.manual_knowledge if d['title'] == '🧠_자가학습'), None)
                    new_mem = f"\n[Q&A]\nQ: {mp}\nA: {reply}\n"
                    if auto_doc: auto_doc['content'] = (auto_doc['content'] + new_mem)[-40000:]
                    else: st.session_state.manual_knowledge.append({"title": "🧠_자가학습", "content": new_mem})
                    save_json(MANUAL_KNOWLEDGE_FILE, st.session_state.manual_knowledge)
                except Exception as e: st.error(e)

# 탭 5: 코어워크(PPT/Canva)
with tab5:
    c_menu = st.radio("메뉴", ["📊 AI 자동 발표 슬라이드 생성", "🎨 Canva 보고서 초안 생성기"], horizontal=True, label_visibility="collapsed")
    if "PPT" in c_menu:
        ppt_t = st.text_input("주제", "생산성 향상 방안")
        if st.button("🚀 PPT 생성"):
            with st.spinner("Generating PPT..."):
                try:
                    p = f'2장짜리 슬라이드 내용 작성. JSON 포맷 필수: {{"s1_title":"제목", "s1_sub":"부제목", "s2_title":"본문제목", "s2_content":"요약내용", "s2_img_prompt":"DALL-E 프롬프트 영문"}} 주제: {ppt_t}'
                    res = client.chat.completions.create(model="gpt-4o", messages=[{"role": "system", "content": "Output valid JSON"}, {"role": "user", "content": p}], response_format={"type": "json_object"})
                    d = json.loads(res.choices[0].message.content)
                    
                    img_url = client.images.generate(model="dall-e-3", prompt=d["s2_img_prompt"], size="1024x1024").data[0].url
                    img_bytes = urllib.request.urlopen(urllib.request.Request(img_url, headers={'User-Agent': 'Mozilla/5.0'})).read()
                    
                    prs = Presentation()
                    sl1 = prs.slides.add_slide(prs.slide_layouts[0])
                    sl1.shapes.title.text, sl1.placeholders[1].text = d.get("s1_title",""), d.get("s1_sub","")
                    sl2 = prs.slides.add_slide(prs.slide_layouts[5])
                    sl2.shapes.title.text = d.get("s2_title","")
                    sl2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)).text_frame.text = d.get("s2_content","")
                    sl2.shapes.add_picture(BytesIO(img_bytes), Inches(5.2), Inches(1.5), width=Inches(4.5))
                    
                    out = BytesIO()
                    prs.save(out)
                    out.seek(0)
                    st.download_button("📥 다운로드 (.pptx)", out, f"{ppt_t}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                    st.image(img_url)
                except Exception as e: st.error(e)
    else:
        cv_t = st.text_input("주제", "종합 분석 보고서")
        cv_d = st.checkbox("데이터 반영", True)
        if st.button("🚀 Canva 텍스트 생성"):
            with st.spinner("Generating..."):
                ds = f"\n[Data]:\n{pd.DataFrame(st.session_state.production_history).tail(10).to_string()}" if cv_d and st.session_state.production_history else ""
                res = client.chat.completions.create(model="gpt-4o", messages=[{"role": "system", "content": "Canva에 붙여넣기 좋은 구조화된 텍스트 작성."}, {"role": "user", "content": f"주제: {cv_t}{ds}"}])
                st.text_area("결과", res.choices[0].message.content, height=300)

# 탭 6: 데이터분석
with tab6:
    a_file = st.file_uploader("분석용 엑셀 업로드", type=["csv", "xlsx"])
    if a_file:
        df_a = pd.read_csv(a_file) if a_file.name.endswith('.csv') else pd.read_excel(a_file)
        with st.expander("미리보기"): st.dataframe(df_a.head())
        for msg in st.session_state.data_chat:
            with st.chat_message(msg["role"]): st.markdown(msg["content"])
        if dp := st.chat_input("데이터 질문"):
            st.session_state.data_chat.append({"role": "user", "content": dp})
            with st.chat_message("user"): st.markdown(dp)
            with st.chat_message("assistant"):
                with st.spinner("Analyzing..."):
                    ds = df_a.to_string()
                    if len(ds) > 20000: ds = f"통계요약:\n{df_a.describe().to_string()}\n상위100행:\n{df_a.head(100).to_string()}"
                    try:
                        res = client.chat.completions.create(model="gpt-4o", messages=[{"role": "system", "content": "데이터 분석가입니다."}] + st.session_state.data_chat[:-1] + [{"role": "user", "content": f"{dp}\n\n[Data]\n{ds}"}], temperature=0.2)
                        st.markdown(res.choices[0].message.content)
                        st.session_state.data_chat.append({"role": "assistant", "content": res.choices[0].message.content})
                    except Exception as e: st.error(e)
